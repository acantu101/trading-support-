#!/usr/bin/env python3
"""
Export all slide decks from challenge-lab.html to individual .pptx files.
Outputs: pptx_export/<deck>.pptx  +  trading-support-slides.zip
"""

import os, sys, zipfile
from pathlib import Path
from bs4 import BeautifulSoup
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Colors ────────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x1A, 0x1A, 0x2E)   # navy background
C_SURF     = RGBColor(0x22, 0x22, 0x38)   # card / panel
C_BORDER   = RGBColor(0x2E, 0x2E, 0x4A)   # border
C_ACCENT   = RGBColor(0xE8, 0x6B, 0x1A)   # orange
C_TEXT     = RGBColor(0xE8, 0xEA, 0xF6)   # main text
C_MUTED    = RGBColor(0x64, 0x74, 0x8B)   # muted / labels
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_CODE_BG  = RGBColor(0x12, 0x12, 0x1E)
C_CODE_TXT = RGBColor(0x52, 0xC8, 0x7A)
C_TIP_BG   = RGBColor(0xFD, 0xEA, 0xD6)
C_TIP_TXT  = RGBColor(0x6A, 0x30, 0x00)
C_NOTE_BG  = RGBColor(0xDD, 0xF0, 0xF0)
C_NOTE_TXT = RGBColor(0x00, 0x4F, 0x4F)
C_WARN_BG  = RGBColor(0xFE, 0xF3, 0xCC)
C_WARN_TXT = RGBColor(0x5A, 0x40, 0x00)

W = Inches(13.33)
H = Inches(7.50)
M = Inches(0.30)   # margin

DECK_NAMES = {
    'fix':        'Market Protocols',
    'kafka':      'Kafka',
    'k8s':        'Kubernetes & Argo',
    'marketdata': 'Market Data',
    'sql':        'SQL',
    'airflow':    'Airflow',
    'linux':      'Linux & Systems',
    'python':     'Python & Bash',
    'aws':        'AWS',
    'networking': 'Networking',
    'git':        'Git',
    'support':    'Support & Incidents',
    'interview':  'Interview Prep',
}


# ── Low-level helpers ─────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    try:
        shape.line.fill.background()
    except Exception:
        pass


def solid_fill_border(shape, fill_color, border_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color


def txb(slide, text, left, top, width, height,
        font='Calibri', size=12, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True):
    """Add a simple single-run text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def multiline_tb(slide, lines, left, top, width, height,
                 font='Calibri', size=12, color=None, bold=False, wrap=True):
    """Add a text box with multiple lines."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    col = color or C_TEXT
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = str(line)
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = col
    return tb


def rect(slide, left, top, width, height, fill=None, border=None):
    """Add a rectangle shape."""
    shp = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
    else:
        try:
            shp.line.fill.background()
        except Exception:
            pass
    return shp


def bullets_tb(slide, items, left, top, width, height, size=10):
    """Add a bullet list text box (▸ marker in orange, text in C_TEXT)."""
    if not items:
        return
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # orange bullet marker
        r1 = p.add_run()
        r1.text = '▸ '
        r1.font.name = 'Calibri'
        r1.font.size = Pt(size)
        r1.font.color.rgb = C_ACCENT
        # text
        r2 = p.add_run()
        r2.text = str(item)
        r2.font.name = 'Calibri'
        r2.font.size = Pt(size)
        r2.font.color.rgb = C_TEXT
    return tb


# ── Content extractors ────────────────────────────────────────────────────────

def get_bullets(ul_el):
    return [li.get_text(' ', strip=True) for li in ul_el.find_all('li', recursive=False)]


def get_title(slide_el):
    """Return (title, subtitle) strings."""
    tl = slide_el.find(class_='sl-title')
    if tl:
        sub = slide_el.find(class_='sl-sub')
        return tl.get_text(' ', strip=True), (sub.get_text(' ', strip=True) if sub else '')
    h2 = slide_el.find('h2')
    if h2:
        p = slide_el.find('p')
        return h2.get_text(' ', strip=True), (p.get_text(' ', strip=True) if p else '')
    return 'Untitled', ''


def content_elements(slide_el):
    """Yield content child elements, skipping header boilerplate."""
    skip = {'sl-num', 'sl-title', 'sl-sub', 'sl-rule'}
    first_p = slide_el.find('p')      # first <p> used as subtitle in new decks
    for child in slide_el.children:
        if not getattr(child, 'name', None):
            continue
        cls = set(child.get('class') or [])
        if cls & skip:
            continue
        if child.name in ('h1', 'h2') and not cls:
            continue
        if child.name == 'p' and not cls and child is first_p:
            continue
        yield child


# ── Slide title bar ───────────────────────────────────────────────────────────

def draw_title_bar(slide, title, subtitle, deck_name):
    """Draw orange left bar + title + subtitle. Returns content_top Y."""
    # Left accent bar
    rect(slide, Inches(0), Inches(0), Inches(0.06), H, fill=C_ACCENT)

    # Deck badge top-right
    badge_w = Inches(2.2)
    txb(slide, deck_name.upper(),
        W - badge_w - Inches(0.1), Inches(0.12), badge_w, Inches(0.28),
        font='Consolas', size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)

    # Title
    txb(slide, title,
        Inches(0.22), Inches(0.12), W - Inches(2.6), Inches(0.70),
        font='Trebuchet MS', size=26, bold=True, color=C_WHITE)

    # Subtitle
    if subtitle:
        txb(slide, subtitle,
            Inches(0.22), Inches(0.80), W - Inches(2.6), Inches(0.34),
            font='Calibri', size=12, color=C_ACCENT)

    # Thin orange divider
    rect(slide, Inches(0.22), Inches(1.14), W - Inches(0.44), Pt(2), fill=C_ACCENT)

    return Inches(1.24)   # content starts here


# ── Content renderers ─────────────────────────────────────────────────────────

def render_card(slide, el, left, top, width, height):
    rect(slide, left, top, width, height, fill=C_SURF, border=C_BORDER)
    cy = top + Inches(0.12)
    h3 = el.find('h3')
    if h3:
        txb(slide, h3.get_text(' ', strip=True).upper(),
            left + Inches(0.14), cy, width - Inches(0.28), Inches(0.28),
            font='Consolas', size=9, bold=True, color=C_ACCENT)
        cy += Inches(0.30)

    # Bullets first
    ul = el.find('ul')
    if ul:
        bullets_tb(slide, get_bullets(ul),
                   left + Inches(0.14), cy,
                   width - Inches(0.28), top + height - cy - Inches(0.1),
                   size=10)
    else:
        # Check for nested pre
        pre = el.find('pre')
        if pre:
            render_code(slide, pre,
                        left + Inches(0.14), cy,
                        width - Inches(0.28), top + height - cy - Inches(0.1))
        else:
            text = el.get_text(' ', strip=True)
            if h3:
                text = text.replace(h3.get_text(' ', strip=True), '', 1).strip()
            if text:
                txb(slide, text,
                    left + Inches(0.14), cy,
                    width - Inches(0.28), top + height - cy - Inches(0.1),
                    font='Calibri', size=10, color=C_TEXT, wrap=True)


def render_code(slide, pre_el, left, top, width, height):
    code = pre_el.get_text('\n')
    lines = code.split('\n')
    if len(lines) > 22:
        lines = lines[:22] + ['…']
    code = '\n'.join(lines)
    rect(slide, left, top, width, height, fill=C_CODE_BG, border=C_BORDER)
    txb(slide, code,
        left + Inches(0.12), top + Inches(0.08),
        width - Inches(0.24), height - Inches(0.18),
        font='Consolas', size=9, color=C_CODE_TXT, wrap=True)


def render_callout(slide, el, left, top, width, height):
    cls = el.get('class') or []
    if 'tip' in cls:
        bg, tc = C_TIP_BG, C_TIP_TXT
    elif 'note' in cls:
        bg, tc = C_NOTE_BG, C_NOTE_TXT
    else:
        bg, tc = C_WARN_BG, C_WARN_TXT
    rect(slide, left, top, width, height, fill=bg, border=RGBColor(0xCC, 0xAA, 0x80))
    txb(slide, el.get_text(' ', strip=True),
        left + Inches(0.14), top + Inches(0.07),
        width - Inches(0.28), height - Inches(0.16),
        font='Calibri', size=10, color=tc, wrap=True)


def render_stats(slide, stats_el, left, top, width, height):
    items = stats_el.find_all(class_='stat')
    n = len(items)
    if n == 0:
        return
    sw = (width - Inches(0.1) * (n - 1)) // n
    for i, stat in enumerate(items):
        sl = left + i * (sw + Inches(0.1))
        rect(slide, sl, top, sw, height, fill=C_SURF, border=C_BORDER)
        val = stat.find(class_='stat-val')
        lbl = stat.find(class_='stat-lbl')
        if val:
            txb(slide, val.get_text(strip=True),
                sl, top + Inches(0.3), sw, Inches(0.7),
                font='Consolas', size=20, bold=True,
                color=C_ACCENT, align=PP_ALIGN.CENTER)
        if lbl:
            txb(slide, lbl.get_text(strip=True).upper(),
                sl, top + Inches(1.0), sw, Inches(0.3),
                font='Calibri', size=9, color=C_MUTED, align=PP_ALIGN.CENTER)


def render_flow(slide, flow_el, left, top, width, height):
    boxes = [el.get_text(' ', strip=True)
             for el in flow_el.find_all(class_='flow-box')]
    text = '  →  '.join(boxes)
    txb(slide, text, left, top, width, height,
        font='Consolas', size=11, color=C_ACCENT, wrap=True)


def render_table(slide, table_el, left, top, width, height):
    rows = table_el.find_all('tr')
    if not rows:
        return
    headers = rows[0].find_all(['th', 'td'])
    ncols = len(headers)
    nrows = len(rows)
    if ncols == 0 or nrows == 0:
        return
    row_h = min(Inches(0.38), height // nrows)
    try:
        tbl_shape = slide.shapes.add_table(nrows, ncols, left, top,
                                            width, row_h * nrows)
        tbl = tbl_shape.table
        for ri, row in enumerate(rows):
            cells = row.find_all(['th', 'td'])
            for ci, cell in enumerate(cells[:ncols]):
                tc = tbl.cell(ri, ci)
                tc.text = cell.get_text(' ', strip=True)
                for para in tc.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(9)
                        run.font.name = 'Calibri'
                        is_header = (ri == 0 or cell.name == 'th')
                        run.font.bold = is_header
                        run.font.color.rgb = C_ACCENT if is_header else C_TEXT
                # Cell fill via XML
                tc_xml = tc._tc
                tcPr = tc_xml.get_or_add_tcPr()
                sf = etree.SubElement(tcPr, qn('a:solidFill'))
                clr = etree.SubElement(sf, qn('a:srgbClr'))
                clr.set('val', '222238' if (ri == 0 or ri % 2 == 0) else '1A1A2E')
    except Exception as e:
        # Fallback: plain text
        text = table_el.get_text(' ', strip=True)[:600]
        txb(slide, text, left, top, width, height,
            font='Calibri', size=10, color=C_TEXT, wrap=True)


def render_content_area(slide, els, left, top, width, height):
    """Lay out all content elements in the available area."""
    # Separate callouts from main content
    callouts = [e for e in els if any(c in (e.get('class') or [])
                                       for c in ('tip', 'note', 'warn'))]
    mains    = [e for e in els if e not in callouts]

    callout_h = Inches(0.52) * len(callouts) + (Inches(0.06) * max(0, len(callouts) - 1))
    main_h = height - (callout_h + Inches(0.08) if callouts else 0)

    # ── Main content ──────────────────────────────────────────────────────
    for el in mains:
        cls = set(el.get('class') or [])
        tag = el.name

        if 'cols' in cls:
            children = [c for c in el.children
                        if getattr(c, 'name', None)]
            n = len(children) or 1
            col_w = (width - Inches(0.1) * (n - 1)) // n
            for i, child in enumerate(children):
                cl = left + i * (col_w + Inches(0.1))
                render_card(slide, child, cl, top, col_w, main_h)

        elif tag == 'pre':
            render_code(slide, el, left, top, width, main_h)

        elif tag == 'table':
            render_table(slide, el, left, top, width, main_h)

        elif 'stats' in cls:
            render_stats(slide, el, left, top, width, main_h)

        elif 'flow' in cls:
            render_flow(slide, el, left, top, width, main_h)

        elif tag == 'ul':
            bullets_tb(slide, get_bullets(el),
                       left, top, width, main_h, size=11)

        elif tag in ('p', 'div'):
            text = el.get_text(' ', strip=True)
            if text:
                txb(slide, text, left, top, width, main_h,
                    font='Calibri', size=11, color=C_TEXT, wrap=True)

    # ── Callouts at bottom ────────────────────────────────────────────────
    cy = top + main_h + Inches(0.06)
    ch = Inches(0.50)
    for callout in callouts:
        render_callout(slide, callout, left, cy, width, ch)
        cy += ch + Inches(0.04)


# ── Main generator ────────────────────────────────────────────────────────────

def build_deck_pptx(deck_id, deck_el, output_dir):
    deck_name = DECK_NAMES.get(deck_id, deck_id.upper())
    wrap = deck_el.find('div', id=f'slides-{deck_id}')
    if not wrap:
        return None
    slide_els = wrap.find_all('div', class_='slide', recursive=False)
    if not slide_els:
        return None

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    for slide_el in slide_els:
        pslide = prs.slides.add_slide(blank)
        set_bg(pslide, C_BG)

        title, subtitle = get_title(slide_el)
        content_top = draw_title_bar(pslide, title, subtitle, deck_name)

        els = list(content_elements(slide_el))
        render_content_area(
            pslide, els,
            left=Inches(0.22),
            top=content_top,
            width=W - Inches(0.44),
            height=H - content_top - Inches(0.12)
        )

    safe_name = deck_name.replace(' ', '_').replace('&', 'and').replace('/', '-')
    fname = f'{deck_id}_{safe_name}.pptx'
    out = os.path.join(output_dir, fname)
    prs.save(out)
    print(f'  OK {fname}  ({len(slide_els)} slides)')
    return out


def main():
    here = Path(__file__).parent.parent
    html_path  = here / 'challenge-lab.html'
    output_dir = here / 'pptx_export'
    zip_path   = here / 'trading-support-slides.zip'

    output_dir.mkdir(exist_ok=True)

    print(f'Parsing {html_path.name} ...')
    with open(html_path, encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')

    saved = []
    for deck_el in soup.find_all('div', class_='deck'):
        deck_id = (deck_el.get('id') or '').replace('deck-', '')
        if not deck_id:
            continue
        print(f'Building: {deck_id}')
        path = build_deck_pptx(deck_id, deck_el, str(output_dir))
        if path:
            saved.append(path)

    print(f'\nZipping {len(saved)} files …')
    with zipfile.ZipFile(str(zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
        for p in saved:
            zf.write(p, Path(p).name)

    size_mb = zip_path.stat().st_size / 1_048_576
    print(f'Done -> {zip_path.name}  ({size_mb:.1f} MB)')


if __name__ == '__main__':
    main()
